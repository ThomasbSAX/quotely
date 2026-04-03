"""
Document ingestion pipeline: PDF / TEX / DOCX / TXT / PPTX / XLSX / CSV / Images → VectorStore

Extraction layers:
  1. .tex  → parse LaTeX source directly (perfect math, zero loss)
  2. .pdf  → pymupdf4llm (markdown, multi-column, OCR fallback) + embedded-image OCR
  3. others → dedicated parsers
"""
from __future__ import annotations
import csv
import hashlib
import io
import json
import os
import re
import shutil
import subprocess
import tempfile
import unicodedata
from pathlib import Path

# Force PyTorch to CPU — avoids SIGSEGV from MPS/Metal init on macOS 26 beta
os.environ.setdefault("PYTORCH_ENABLE_MPS_FALLBACK", "1")
os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")

from vector_store import VectorStore
import fitz  # PyMuPDF
import pymupdf4llm
from docx import Document as DocxDocument

# Optional OCR dependencies
try:
    import pytesseract
    from PIL import Image as PILImage
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# Optional PowerPoint / Excel
try:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN  # noqa: F401
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

# Shared singletons (imported by other modules)
DB_PATH     = Path(__file__).parent.parent / "data" / "db"
DATA_PATH   = Path(__file__).parent.parent / "data"        # watched root (all subdirs)
PAPERS_PATH = DATA_PATH / "papers"                          # kept for backwards compat

# LibreOffice path for legacy formats (.doc, .ppt, .odt, .rtf)
_SOFFICE = shutil.which("soffice") or "/opt/homebrew/bin/soffice"
LIBREOFFICE_AVAILABLE = Path(_SOFFICE).exists() if _SOFFICE else False

_collection: VectorStore | None = None
_model = None  # SentenceTransformer, loaded lazily on first use

CHUNK_SIZE = 300   # tokens approx (we use words as proxy)
CHUNK_OVERLAP = 50


def get_collection() -> VectorStore:
    global _collection
    if _collection is None:
        _collection = VectorStore(DB_PATH / "vectors.pkl")
    return _collection


def get_model():
    global _model
    if _model is None:
        from sentence_transformers import SentenceTransformer  # lazy — avoids PyTorch crash on startup
        print("[RagCite] Loading embedding model (first time ~15s)...")
        _model = SentenceTransformer("all-MiniLM-L6-v2", device="cpu")
        print("[RagCite] Model ready.")
    return _model


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def _extract_pdf(path: Path) -> tuple[str, dict]:
    """
    Primary: pymupdf4llm → markdown with structure + OCR fallback on complex pages.
    Secondary: OCR of embedded images (figures, diagrams).
    """
    doc = fitz.open(str(path))
    meta = doc.metadata or {}

    # --- pymupdf4llm: markdown with multi-column layout, table support, OCR ---
    try:
        md_text = pymupdf4llm.to_markdown(str(path), show_progress=False)
    except Exception:
        # Fallback to raw text if pymupdf4llm fails
        md_text = "\n".join(page.get_text("text") for page in doc)

    # --- OCR of embedded images (figures, diagrams, formula images) ---
    img_texts: list[str] = []
    if OCR_AVAILABLE:
        seen_xrefs: set[int] = set()
        for page in doc:
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)
                try:
                    base = doc.extract_image(xref)
                    w, h = base["width"], base["height"]
                    # Skip tiny images (icons, decorations < 50×50 px)
                    if w < 50 or h < 50:
                        continue
                    pil_img = PILImage.open(io.BytesIO(base["image"]))
                    ocr = pytesseract.image_to_string(pil_img, lang="fra+eng").strip()
                    if ocr:
                        img_texts.append(f"[Figure/Image]: {ocr}")
                except Exception:
                    pass

    doc.close()

    full_text = md_text
    if img_texts:
        full_text += "\n\n" + "\n\n".join(img_texts)

    # Fix LaTeX font encoding artifacts across the full text
    full_text = _fix_latex_encoding(full_text)

    title = meta.get("title", "").strip()
    author = meta.get("author", "").strip()
    year = _extract_year(meta.get("creationDate", "") + " " + full_text[:2000])

    if not title:
        # Look for first markdown heading or substantial line (skip date-only lines)
        date_re = re.compile(r"^\d{1,2}\s+\w+\s+\d{4}$|^\d{4}$")
        for line in full_text.splitlines():
            line = line.lstrip("#").strip()
            if len(line) > 15 and not date_re.match(line):
                title = line[:120]
                break

    # Fix common LaTeX font encoding artifacts in extracted text
    title = _fix_latex_encoding(title)
    author = _fix_latex_encoding(author)

    return full_text, {"title": title, "authors": author, "year": year}


def _extract_tex(path: Path) -> tuple[str, dict]:
    """
    Parse LaTeX source file directly — preserves math perfectly.
    Strips LaTeX commands but keeps: section bodies, $...$ inline math,
    \\begin{equation}...\\end{equation}, \\begin{align}...
    """
    raw = path.read_text(encoding="utf-8", errors="ignore")

    # Extract metadata from preamble
    def _latex_arg(cmd: str) -> str:
        m = re.search(rf"\\{cmd}\{{([^}}]{{1,300}})\}}", raw)
        return m.group(1).strip() if m else ""

    title = _latex_arg("title") or path.stem
    authors = _latex_arg("author")
    year = _extract_year(_latex_arg("date") + raw[:3000])

    # Remove comments
    src = re.sub(r"%.*$", "", raw, flags=re.MULTILINE)

    # Keep only document body
    body_m = re.search(r"\\begin\{document\}(.*?)\\end\{document\}", src, re.DOTALL)
    body = body_m.group(1) if body_m else src

    # Replace section commands with readable headers
    body = re.sub(r"\\(?:sub)*section\*?\{([^}]+)\}", r"\n\n== \1 ==\n", body)

    # Preserve math blocks as-is (they're the key content)
    # Just clean surrounding whitespace, don't strip them

    # Strip remaining LaTeX commands but keep their arguments when meaningful
    # Remove environments that aren't content (figure env wrapper, but keep caption)
    body = re.sub(r"\\caption\{([^}]+)\}", r"[Caption: \1]", body)
    body = re.sub(r"\\label\{[^}]+\}", "", body)
    body = re.sub(r"\\ref\{[^}]+\}", "[ref]", body)
    body = re.sub(r"\\cite\{[^}]+\}", "[cite]", body)
    # Keep text formatting content
    body = re.sub(r"\\(?:textbf|textit|emph|underline|text)\{([^}]+)\}", r"\1", body)
    # Remove remaining single-argument commands (figures, includes, etc.)
    body = re.sub(r"\\[a-zA-Z]+\*?\{[^}]*\}", "", body)
    # Remove standalone commands
    body = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])?", " ", body)
    # Clean up excess whitespace
    body = re.sub(r"\n{3,}", "\n\n", body)

    return body.strip(), {"title": title, "authors": authors, "year": year}


def _extract_docx(path: Path) -> tuple[str, dict]:
    doc = DocxDocument(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    full_text = "\n".join(parts)
    title = parts[0][:120] if parts else path.stem
    year = _extract_year(full_text[:500])
    return full_text, {"title": title, "authors": "", "year": year}


def _extract_txt(path: Path) -> tuple[str, dict]:
    full_text = path.read_text(encoding="utf-8", errors="ignore").lstrip("\ufeff")
    lines = full_text.splitlines()
    title = next((l.strip()[:120] for l in lines if l.strip()), path.stem)
    year = _extract_year(full_text[:500])
    return full_text, {"title": title, "authors": "", "year": year}


def _extract_pptx(path: Path) -> tuple[str, dict]:
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")
    prs = Presentation(str(path))
    title = path.stem
    # Try to get title from first slide's title placeholder
    if prs.slides:
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PLACEHOLDER
                t = shape.text.strip()
                if t:
                    title = t[:120]
                    break
        if title == path.stem:
            # Fallback: first non-empty text in slide 1
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame:
                    t = shape.text.strip()
                    if len(t) > 5:
                        title = t[:120]
                        break
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    full_text = "\n".join(texts)
    year = _extract_year(full_text[:5000])
    return full_text, {"title": title, "authors": "", "year": year}


def _extract_xlsx(path: Path) -> tuple[str, dict]:
    if not XLSX_AVAILABLE:
        raise ImportError("openpyxl not installed. Run: pip install openpyxl")
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    rows = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
    wb.close()
    full_text = "\n".join(rows)
    year = _extract_year(full_text[:5000])
    return full_text, {"title": path.stem, "authors": "", "year": year}


def _extract_csv(path: Path) -> tuple[str, dict]:
    rows = []
    with open(path, encoding="utf-8", errors="ignore", newline="") as f:
        reader = csv.reader(f)
        for row in reader:
            if any(cell.strip() for cell in row):
                rows.append(" | ".join(row))
    full_text = "\n".join(rows)
    year = _extract_year(full_text[:5000])
    return full_text, {"title": path.stem, "authors": "", "year": year}


def _extract_image(path: Path) -> tuple[str, dict]:
    if not OCR_AVAILABLE:
        raise ImportError(
            "pytesseract/Pillow not installed or Tesseract absent. "
            "Run: brew install tesseract && pip install pytesseract Pillow"
        )
    img = PILImage.open(str(path))
    full_text = pytesseract.image_to_string(img, lang="fra+eng")
    year = _extract_year(full_text[:5000])
    lines = full_text.splitlines()
    title = next((l.strip()[:120] for l in lines if len(l.strip()) > 5), path.stem)
    return full_text, {"title": title, "authors": "", "year": year}


def _extract_doc_via_libreoffice(path: Path) -> tuple[str, dict]:
    """Convert .doc / .ppt / .odt / .rtf to text via LibreOffice headless."""
    if not LIBREOFFICE_AVAILABLE:
        raise ImportError(
            "LibreOffice not found. Install with: brew install --cask libreoffice"
        )
    with tempfile.TemporaryDirectory() as tmpdir:
        subprocess.run(
            [_SOFFICE, "--headless", "--convert-to", "txt:Text",
             "--outdir", tmpdir, str(path)],
            capture_output=True,
            timeout=90,
            check=True,
        )
        txt_path = Path(tmpdir) / (path.stem + ".txt")
        if not txt_path.exists():
            raise RuntimeError(f"LibreOffice produced no output for {path.name}")
        return _extract_txt(txt_path)


def _extract_ipynb(path: Path) -> tuple[str, dict]:
    """Extract text from a Jupyter notebook (.ipynb)."""
    nb = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    parts: list[str] = []
    for cell in nb.get("cells", []):
        src = "".join(cell.get("source", []))
        if src.strip():
            parts.append(src)
        for output in cell.get("outputs", []):
            txt = "".join(output.get("text", []))
            if txt.strip():
                parts.append(txt)
    full_text = "\n\n".join(parts)
    title = next(
        (l.lstrip("# ").strip()[:120] for l in full_text.splitlines()
         if l.strip() and not l.startswith("```")),
        path.stem,
    )
    year = _extract_year(full_text[:5000])
    return full_text, {"title": title, "authors": "", "year": year}


def _fix_latex_encoding(text: str) -> str:
    """Fix common LaTeX font encoding artifacts in PyMuPDF text extraction."""
    # Combining accent + letter → precomposed Unicode character
    replacements = [
        ("´a", "á"), ("´e", "é"), ("´i", "í"), ("´o", "ó"), ("´u", "ú"),
        ("`a", "à"), ("`e", "è"), ("`i", "ì"), ("`o", "ò"), ("`u", "ù"),
        ("ˆa", "â"), ("ˆe", "ê"), ("ˆi", "î"), ("ˆo", "ô"), ("ˆu", "û"),
        ("¨a", "ä"), ("¨e", "ë"), ("¨i", "ï"), ("¨o", "ö"), ("¨u", "ü"),
        ("c¸", "ç"), ("c,", "ç"),
        # Common ligature issues
        ("ﬁ", "fi"), ("ﬂ", "fl"), ("ﬀ", "ff"), ("ﬃ", "ffi"), ("ﬄ", "ffl"),
    ]
    for bad, good in replacements:
        text = text.replace(bad, good)
    return text


def _extract_year(text: str) -> str:
    m = re.search(r"\b(?:19|20)\d{2}\b", text)
    return m.group(0) if m else "n.d."


def _infer_author_from_path(path: Path) -> str:
    """
    Guess author from folder path by extracting capitalized proper nouns
    from parent directory names.
    e.g. 'EcritspublicationsHarbulot' → 'Harbulot'
    """
    SKIP_LOWER = {
        'data', 'papers', 'documents', 'db', 'downloads', 'users',
        '.venv', 'backend', 'extension', 'node_modules', 'ragflow',
        'projets', 'perso', 'library', 'desktop',
    }
    GENERIC = {
        'Ecrits', 'Publications', 'Articles', 'Documents', 'Markdowns',
        'Projets', 'Downloads', 'Users', 'Library', 'Desktop', 'Papers',
        'Notes', 'Files', 'Folder',
    }
    for part in reversed(path.parent.parts):
        if part.lower() in SKIP_LOWER or part.startswith('.'):
            continue
        words = re.findall(r'[A-Z][a-z]{3,}', part)
        candidates = [w for w in words if w not in GENERIC]
        if candidates:
            return candidates[-1]
    return ""


def _infer_author_from_text(text: str) -> str:
    """Look for common author attribution patterns in the first 600 chars."""
    patterns = [
        r'[Pp]ar\s+([A-ZÀ-Ü][a-zà-ü]+(?:\s+[A-ZÀ-Ü][a-zà-ü]+)?)',
        r'[Aa]uteur\s*[:\-]\s*([A-ZÀ-Ü][a-zà-ü]+(?:\s+[A-ZÀ-Ü][a-zà-ü]+)?)',
        r'[Aa]uthor\s*[:\-]\s*([A-ZÀ-Ü][a-zà-ü]+(?:\s+[A-ZÀ-Ü][a-zà-ü]+)?)',
    ]
    snippet = text[:600]
    for pat in patterns:
        m = re.search(pat, snippet)
        if m:
            return m.group(1).strip()
    return ""


# ---------------------------------------------------------------------------
# BibTeX helpers
# ---------------------------------------------------------------------------

def _slugify(text: str) -> str:
    text = unicodedata.normalize("NFD", text)
    text = "".join(c for c in text if unicodedata.category(c) != "Mn")
    text = re.sub(r"[^a-zA-Z0-9]", "", text)
    return text


def _extract_lastname(author_str: str) -> str:
    """Extract a clean last name from author string (name or email)."""
    # Take first author only
    first = re.split(r"[;]", author_str)[0].strip()
    # Email: louise.budzynski@psl.eu → budzynski
    if "@" in first:
        local = first.split("@")[0]           # louise.budzynski
        parts = re.split(r"[.\-_]", local)
        last = parts[-1] if len(parts) > 1 else parts[0]
        return _slugify(last).capitalize()
    # "Last, First" format
    if "," in first:
        return _slugify(first.split(",")[0].strip()).capitalize()
    # "First Last" or "First M. Last"
    parts = first.split()
    return _slugify(parts[-1]).capitalize() if parts else "Unknown"


def _make_bibtex_key(authors: str, year: str, title: str) -> str:
    lastname = _extract_lastname(authors) if authors else "Unknown"
    # Short year: 2025 → 25, n.d. → nd
    if re.match(r"^\d{4}$", year):
        short_year = year[-2:]
    elif year == "n.d.":
        short_year = "nd"
    else:
        short_year = year
    # One distinctive title word (skip stop words + markdown artifacts)
    stop = {"a", "an", "the", "de", "du", "des", "le", "la", "les", "et",
            "en", "un", "une", "of", "in", "on", "to", "for", "cours", "n"}
    clean_title = re.sub(r"[*#\-–—]+", " ", title)  # strip markdown bold/headers
    words = [_slugify(w) for w in clean_title.split()
             if _slugify(w) and _slugify(w).lower() not in stop and len(_slugify(w)) > 2]
    title_word = words[0].capitalize() if words else "Untitled"
    return f"{lastname}{short_year}{title_word}"


def _make_bibtex_entry(key: str, title: str, authors: str, year: str, file_path: str) -> str:
    return (
        f"@article{{{key},\n"
        f"  title   = {{{title}}},\n"
        f"  author  = {{{authors or 'Unknown'}}},\n"
        f"  year    = {{{year}}},\n"
        f"  note    = {{File: {file_path}}}\n"
        f"}}"
    )


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def _chunk_text(text: str) -> list[str]:
    """Fallback: word-count chunking with overlap."""
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + CHUNK_SIZE
        chunk = " ".join(words[start:end])
        if chunk.strip():
            chunks.append(chunk)
        start = end - CHUNK_OVERLAP
        if start >= len(words):
            break
    return chunks


def _chunk_markdown(text: str) -> list[str]:
    """
    Section-aware chunking for markdown/structured text.
    Splits on ## / ### headers, keeping each section as a unit.
    Large sections are further split with overlap. Tiny sections are merged.
    Figure captions are kept as standalone chunks for better retrieval.
    """
    # Split on markdown headers (## or ###) or LaTeX-style == Section ==
    section_re = re.compile(r"(?=^#{1,3} |\n#{1,3} |^== .+ ==$|\n== .+ ==$)", re.MULTILINE)
    parts = section_re.split(text)

    chunks: list[str] = []
    buffer = ""

    for part in parts:
        part = part.strip()
        if not part:
            continue

        # Extract figure captions as dedicated chunks
        captions = re.findall(r"\[Caption:[^\]]{10,}\]|\[Figure/Image\]:[^\n]{10,}", part)
        for cap in captions:
            if cap.strip():
                chunks.append(cap.strip())

        words = part.split()
        if len(words) <= 20:
            # Too short — accumulate in buffer
            buffer += " " + part
        elif len(words) <= CHUNK_SIZE:
            combined = (buffer + " " + part).strip()
            if len(combined.split()) > CHUNK_SIZE:
                if buffer.strip():
                    chunks.append(buffer.strip())
                buffer = part
            else:
                buffer = combined
        else:
            # Flush buffer first
            if buffer.strip():
                chunks.append(buffer.strip())
                buffer = ""
            # Split large section with overlap, prefixed with section header
            header_m = re.match(r"(#{1,3} [^\n]+|== [^\n]+ ==)", part)
            header = header_m.group(0) + " — " if header_m else ""
            sub = _chunk_text(part)
            for s in sub:
                chunks.append((header + s).strip())

    if buffer.strip():
        chunks.append(buffer.strip())

    # Final filter: drop chunks that are pure whitespace or very short
    return [c for c in chunks if len(c.split()) >= 8]


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(65536), b""):
            h.update(block)
    return h.hexdigest()


def is_already_indexed(path: Path) -> bool:
    col = get_collection()
    sha = file_hash(path)
    results = col.get(where={"file_hash": sha}, limit=1)
    return len(results["ids"]) > 0


def ingest_file(path: Path) -> dict:
    """Parse, chunk, embed and store a document. Returns summary dict."""
    path = Path(path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        text, meta = _extract_pdf(path)
        chunks = _chunk_markdown(text)       # section-aware for structured markdown
    elif suffix == ".tex":
        text, meta = _extract_tex(path)
        chunks = _chunk_markdown(text)       # LaTeX body is already structured
    elif suffix == ".docx":
        text, meta = _extract_docx(path)
        chunks = _chunk_text(text)
    elif suffix in (".doc", ".odt", ".rtf", ".ppt"):
        text, meta = _extract_doc_via_libreoffice(path)
        chunks = _chunk_text(text)
    elif suffix in (".md", ".txt"):
        text, meta = _extract_txt(path)
        chunks = _chunk_markdown(text)
    elif suffix == ".pptx":
        text, meta = _extract_pptx(path)
        chunks = _chunk_text(text)
    elif suffix in (".xlsx", ".xls"):
        text, meta = _extract_xlsx(path)
        chunks = _chunk_text(text)
    elif suffix == ".csv":
        text, meta = _extract_csv(path)
        chunks = _chunk_text(text)
    elif suffix == ".ipynb":
        text, meta = _extract_ipynb(path)
        chunks = _chunk_markdown(text)
    elif suffix in (".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp"):
        text, meta = _extract_image(path)
        chunks = _chunk_text(text)
    else:
        raise ValueError(f"Unsupported format: {suffix}")

    title = meta["title"] or path.stem
    authors = meta["authors"]
    year = meta["year"]

    # Fill missing author from path structure then from text content
    if not authors:
        authors = _infer_author_from_path(path) or _infer_author_from_text(text)

    bibtex_key = _make_bibtex_key(authors, year, title)
    bibtex_entry = _make_bibtex_entry(bibtex_key, title, authors, year, str(path))
    sha = file_hash(path)

    if not chunks:
        raise ValueError("No text extracted from document.")

    # Prepend a dedicated title chunk for direct title-based retrieval
    title_chunk = title + (f" — {authors}" if authors else "")
    all_chunks = [title_chunk] + chunks

    # Embed with title prefix so every chunk carries the article title →
    # boosts retrieval when the query mentions title-related terms
    embed_texts = [f"Titre: {title}\n{chunk}" for chunk in all_chunks]

    model = get_model()
    embeddings = model.encode(embed_texts, show_progress_bar=False).tolist()

    col = get_collection()
    # Title chunk → _title ID; content chunks → _chunk0, _chunk1, …
    ids = [f"{bibtex_key}_title"] + [f"{bibtex_key}_chunk{i}" for i in range(len(chunks))]
    base_meta = {
        "bibtex_key": bibtex_key,
        "title": title,
        "authors": authors,
        "year": year,
        "file_path": str(path),
        "bibtex_entry": bibtex_entry,
        "file_hash": sha,
    }
    metadatas = [
        {**base_meta, "chunk_index": -1},  # title chunk
        *[{**base_meta, "chunk_index": i} for i in range(len(chunks))],
    ]

    col.upsert(ids=ids, embeddings=embeddings, documents=all_chunks, metadatas=metadatas)

    print(f"[Quotely] Indexed '{title}' ({len(all_chunks)} chunks) → {bibtex_key}")
    return {
        "bibtex_key": bibtex_key,
        "title": title,
        "chunks_indexed": len(all_chunks),
    }
